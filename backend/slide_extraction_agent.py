from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
import io

app = Flask(__name__)
CORS(app)

@app.route("/extract", methods=["POST"])
def extract():
    f = request.files["file"]
    ppt = Presentation(io.BytesIO(f.read()))
    slides = []
    for slide in ppt.slides:
        text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slides.append(text)
    return jsonify({"slides": slides})

if __name__ == "__main__":
    app.run(port=5001, debug=True)
